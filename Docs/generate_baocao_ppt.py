# -*- coding: utf-8 -*-
"""Bao cao tien do QLSV - python Docs/generate_baocao_ppt.py"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent / "BaoCaoTienDo_QLSV.pptx"

# Brand
PRIMARY = RGBColor(31, 30, 68)
PRIMARY_SOFT = RGBColor(45, 44, 90)
ACCENT = RGBColor(60, 179, 113)
ACCENT_DARK = RGBColor(46, 139, 87)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(38, 38, 38)
MUTED = RGBColor(105, 105, 115)
LIGHT_BG = RGBColor(247, 248, 252)
BORDER = RGBColor(218, 222, 232)
CODE_BG = RGBColor(240, 242, 248)

FONT = "Segoe UI"
FONT_TITLE = "Segoe UI"  # bold via flag
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

_footer_idx = 0
_footer_total = 0


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, l, t, w, h, fill, border=False, radius=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(kind, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = BORDER
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh


def _text(slide, l, t, w, h, text, size=16, bold=False, color=TEXT, align=PP_ALIGN.LEFT, font=None):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font or FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def _run(p, text, *, size=15, bold=False, color=TEXT, italic=False):
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return r


def _new_para(tf, first=False):
    return tf.paragraphs[0] if first else tf.add_paragraph()


def _field_list(slide, l, t, w, h, fields, numbered=True):
    """Danh sach dang: 1. Nhan de: Noi dung (nhan dam, mau khac)."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, (label, value) in enumerate(fields):
        p = _new_para(tf, i == 0)
        p.space_after = Pt(16)
        p.line_spacing = 1.25
        p.space_before = Pt(4)
        if numbered:
            _run(p, f"{i + 1}.  ", size=14, bold=True, color=ACCENT_DARK)
        _run(p, f"{label}: ", size=15, bold=True, color=PRIMARY)
        _run(p, value, size=15, bold=False, color=TEXT)
    return box


def _bullets(slide, l, t, w, h, items, size=15, color=TEXT, numbered=False, title_size=None):
    """Bullet thuong hoac (tieu de phu, mo ta) neu item la tuple."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    ts = title_size or (size + 1)

    for i, item in enumerate(items):
        if isinstance(item, tuple):
            title, desc = item[0], item[1] if len(item) > 1 else None
            p = _new_para(tf, i == 0)
            p.space_after = Pt(12)
            p.line_spacing = 1.2
            if numbered:
                _run(p, f"{i + 1}.  ", size=14, bold=True, color=ACCENT_DARK)
            _run(p, title, size=ts, bold=True, color=PRIMARY)
            if desc:
                p2 = tf.add_paragraph()
                p2.level = 1
                p2.space_after = Pt(10)
                indent = "     " if numbered else "   "
                _run(p2, indent + desc, size=size, color=MUTED)
            continue

        p = _new_para(tf, i == 0)
        p.space_after = Pt(10)
        p.line_spacing = 1.2
        if numbered:
            _run(p, f"{i + 1}.  ", size=14, bold=True, color=ACCENT_DARK)
            _run(p, str(item), size=size, color=color)
        else:
            _run(p, "•  ", size=size, bold=True, color=ACCENT)
            _run(p, str(item), size=size, color=color)
    return box


def _footer(slide):
    global _footer_idx
    _footer_idx += 1
    _rect(slide, Inches(0), SLIDE_H - Inches(0.42), SLIDE_W, Inches(0.42), PRIMARY)
    _text(slide, Inches(0.5), SLIDE_H - Inches(0.36), Inches(5), Inches(0.28),
          "TDMU  |  Hệ thống Quản lý Sinh viên", size=10, color=WHITE)
    _text(slide, SLIDE_W - Inches(1.3), SLIDE_H - Inches(0.36), Inches(0.8), Inches(0.28),
          f"{_footer_idx}/{_footer_total}", size=10, color=WHITE, align=PP_ALIGN.RIGHT)


def _header(slide, title, section_tag=None):
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.02), PRIMARY)
    _rect(slide, Inches(0), Inches(1.02), SLIDE_W, Inches(0.055), ACCENT)
    if section_tag:
        _text(slide, Inches(0.5), Inches(0.08), Inches(3), Inches(0.28), section_tag.upper(),
              size=10, bold=True, color=ACCENT)
        _text(slide, Inches(0.5), Inches(0.32), Inches(12), Inches(0.62), title,
              size=26, bold=True, color=WHITE, font=FONT_TITLE)
    else:
        _text(slide, Inches(0.5), Inches(0.24), Inches(12), Inches(0.62), title,
              size=26, bold=True, color=WHITE, font=FONT_TITLE)


def _base_content(slide, title, section=None):
    _rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_BG)
    _header(slide, title, section)
    _footer(slide)


def _image_placeholder(slide, l, t, w, h, label):
    _rect(slide, l, t, w, h, RGBColor(252, 252, 254), border=True, radius=True)
    _text(slide, l, t + h / 2 - Inches(0.35), w, Inches(0.35),
          "[ Chèn hình ảnh ]", size=13, color=MUTED, align=PP_ALIGN.CENTER)
    _text(slide, l, t + h / 2 + Inches(0.05), w, Inches(0.3),
          label, size=11, color=ACCENT_DARK, align=PP_ALIGN.CENTER)


# ---- Slide builders ----

def slide_cover(prs):
    s = _blank(prs)
    _rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, PRIMARY)
    _rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.14), ACCENT)
    _rect(s, Inches(0.7), Inches(1.85), Inches(0.08), Inches(2.8), ACCENT)
    _text(s, Inches(1.05), Inches(1.75), Inches(11), Inches(1.1),
          "HỆ THỐNG QUẢN LÝ SINH VIÊN", size=38, bold=True, color=WHITE, font=FONT_TITLE)
    _text(s, Inches(1.05), Inches(2.85), Inches(11), Inches(0.55),
          "Báo cáo tiến độ dự án cuối kì", size=22, color=ACCENT)
    _text(s, Inches(1.05), Inches(3.55), Inches(11), Inches(0.45),
          "Nhóm MRBEAST  ·  Trường Đại học Thủ Dầu Một  ·  16/05/2026", size=15,
          color=RGBColor(190, 190, 210))
    for i, (a, b) in enumerate([("Client", "WinForms UI"), ("Server", "ASP.NET API"), ("Database", "SQL Server")]):
        x = Inches(1.05 + i * 3.85)
        _rect(s, x, Inches(4.65), Inches(3.45), Inches(1.55), PRIMARY_SOFT, border=True, radius=True)
        _text(s, x + Inches(0.2), Inches(4.85), Inches(3), Inches(0.4), a, size=17, bold=True, color=ACCENT)
        _text(s, x + Inches(0.2), Inches(5.3), Inches(3), Inches(0.35), b, size=13, color=WHITE)


def slide_agenda(prs):
    s = _blank(prs)
    _base_content(s, "Nội dung trình bày", "Phần mở đầu")
    items = [
        ("I", "Giới thiệu đề tài và mục tiêu"),
        ("II", "Cấu trúc thư mục dự án"),
        ("III", "Kiến trúc 3 tầng và pipeline xử lý"),
        ("IV", "Ứng dụng làm gì? (chức năng chính)"),
        ("V", "Tiến độ Client (giao diện)"),
        ("VI", "Tiến độ Server (API)"),
        ("VII", "Cơ sở dữ liệu: bảng, SP, trigger"),
        ("VIII", "Luồng cập nhật điểm và trigger tự động"),
        ("IX", "Demo và kế hoạch tiếp theo"),
    ]
    y = Inches(1.35)
    for num, label in items:
        _rect(s, Inches(0.7), y, Inches(0.55), Inches(0.48), ACCENT, radius=True)
        _text(s, Inches(0.7), y + Inches(0.06), Inches(0.55), Inches(0.38), num,
              size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _text(s, Inches(1.45), y + Inches(0.06), Inches(10.5), Inches(0.4), label, size=17, color=TEXT)
        y += Inches(0.58)


def slide_section(prs, num, title, subtitle=""):
    s = _blank(prs)
    _rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, PRIMARY)
    _rect(s, Inches(0), Inches(3.2), SLIDE_W, Inches(0.08), ACCENT)
    _text(s, Inches(0.8), Inches(2.2), Inches(3), Inches(0.7), num, size=52, bold=True, color=ACCENT, font=FONT_TITLE)
    _text(s, Inches(0.8), Inches(3.45), Inches(11), Inches(1), title, size=34, bold=True, color=WHITE, font=FONT_TITLE)
    if subtitle:
        _text(s, Inches(0.8), Inches(4.45), Inches(11), Inches(0.6), subtitle, size=18, color=RGBColor(200, 200, 220))


def slide_intro(prs):
    s = _blank(prs)
    _base_content(s, "Giới thiệu đề tài", "Phần I")
    _field_list(s, Inches(0.65), Inches(1.35), Inches(6.2), Inches(5.2), [
        ("Tên đề tài", "Quản lý Sinh viên (QLSV)"),
        ("Công nghệ", "C# WinForms + ASP.NET Core + SQL Server"),
        ("Mục tiêu", "Số hóa quản lý SV, môn học, điểm, xếp hạng"),
        ("Phạm vi", "3 tầng Client — Server — Database"),
        ("Thời gian", "Khoảng 2–3 tuần (đang tích hợp)"),
    ])
    _image_placeholder(s, Inches(7.1), Inches(1.35), Inches(5.6), Inches(4.9), "Logo TDMU / ảnh nhóm")


def slide_app_features(prs):
    s = _blank(prs)
    _base_content(s, "Ứng dụng làm gì?", "Phần IV")
    cards = [
        ("Quản lý SV", "Thêm, sửa, tìm kiếm, lọc lớp, xóa hàng loạt"),
        ("Quản lý môn", "CRUD môn học, đăng ký môn cho SV"),
        ("Điểm số", "Nhập QT/GK/CK, xem GPA, học lực"),
        ("Báo cáo", "Xếp hạng, học bổng, cảnh báo học vụ"),
    ]
    for i, (t, d) in enumerate(cards):
        col, row = i % 2, i // 2
        x = Inches(0.65 + col * 6.2)
        y = Inches(1.35 + row * 2.45)
        _rect(s, x, y, Inches(5.85), Inches(2.15), WHITE, border=True, radius=True)
        _rect(s, x, y, Inches(5.85), Inches(0.42), ACCENT, radius=True)
        _text(s, x + Inches(0.2), y + Inches(0.06), Inches(5.4), Inches(0.35), t,
              size=15, bold=True, color=WHITE, font=FONT_TITLE)
        _text(s, x + Inches(0.2), y + Inches(0.55), Inches(5.4), Inches(1.4), d, size=14, color=TEXT)


def slide_folder_tree(prs):
    s = _blank(prs)
    _base_content(s, "Cau truc thu muc du an", "Phan II")
    tree = """QLSV_System/
├── Client/          # WinForms - giao dien
│   ├── View/        # frmLogin, frmMain, frmSv...
│   ├── Controller/  # MainController
│   └── Helper/
├── Server/          # ASP.NET Core API
│   ├── Controllers/ # (dang mo rong)
│   └── Program.cs
├── Database/        # SQL scripts
│   ├── CSDL.sql
│   ├── Procetures/
│   ├── Functions/
│   ├── Triggers/
│   └── Views/
└── Shared/          # Model dung chung
    └── Objects/     # SinhVien, MonHoc..."""
    _rect(s, Inches(0.65), Inches(1.3), Inches(6.4), Inches(5.35), CODE_BG, border=True, radius=True)
    _text(s, Inches(0.85), Inches(1.45), Inches(6), Inches(5.1), tree, size=12, color=TEXT, font="Consolas")
    _field_list(s, Inches(7.3), Inches(1.35), Inches(5.5), Inches(2.8), [
        ("Client", "UI hoàn thiện, chưa gọi API"),
        ("Server", "Khung ASP.NET, chưa Repository"),
        ("Database", "Tập trung logic nghiệp vụ"),
        ("Shared", "Model dùng chung Client/Server"),
    ])
    _image_placeholder(s, Inches(7.3), Inches(4.5), Inches(5.5), Inches(2.1), "Screenshot Solution Explorer")


def slide_pipeline(prs):
    s = _blank(prs)
    _base_content(s, "Pipeline xu ly yeu cau (thiet ke)", "Phan III")
    steps = [
        ("1", "Client", "User bam Luu diem\nfrmChitietDiem"),
        ("2", "HTTP", "POST /api/diem\n(JSON MaSV, MaMH, diem...)"),
        ("3", "Server", "Controller ->\nRepository"),
        ("4", "SQL", "EXEC CapNhatDiem\n@MaSV, @MaMH..."),
        ("5", "Trigger", "trg_DiemThi_\nTinhTongKet"),
        ("6", "DB", "Cap nhat TongKet,\nGPA, HocLuc"),
    ]
    x0 = Inches(0.45)
    w = Inches(1.95)
    gap = Inches(0.22)
    for i, (num, title, desc) in enumerate(steps):
        x = x0 + i * (w + gap)
        y = Inches(2.0)
        _rect(s, x, y, w, Inches(2.85), WHITE, border=True, radius=True)
        _rect(s, x, y, w, Inches(0.38), PRIMARY, radius=True)
        _text(s, x, y + Inches(0.05), w, Inches(0.3), num, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _text(s, x + Inches(0.1), y + Inches(0.48), w - Inches(0.2), Inches(0.4), title,
              size=14, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER, font=FONT_TITLE)
        _text(s, x + Inches(0.1), y + Inches(0.95), w - Inches(0.2), Inches(1.6), desc,
              size=11, color=TEXT, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            ax = x + w + Inches(0.04)
            _text(s, ax, y + Inches(1.2), gap, Inches(0.4), ">", size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    _rect(s, Inches(0.65), Inches(5.15), Inches(12), Inches(0.75), RGBColor(232, 245, 236), radius=True)
    _text(s, Inches(0.85), Inches(5.28), Inches(11.6), Inches(0.55),
          "Hien tai: UI + Database san sang  |  Server dang xay dung lop goi procedure",
          size=13, color=ACCENT_DARK)


def slide_trigger_flow(prs):
    s = _blank(prs)
    _base_content(s, "Khi sua diem: Trigger tu dong tinh", "Phan VIII")
    _bullets(s, Inches(0.65), Inches(1.3), Inches(6.0), Inches(5.5), [
        ("Nhập điểm", "Giảng viên nhập QT, GK, CK trên form Client"),
        ("Gọi procedure", "Server: EXEC CapNhatDiem ..."),
        ("Ghi DB", "INSERT / UPDATE bảng DiemThi"),
        ("Trigger kích hoạt", "AFTER INSERT, UPDATE trên DiemThi"),
        ("Tính điểm môn", "fn_TinhDiemTongKet (20-30-50%) → TongKet"),
        ("Tính GPA", "Trung bình có trọng số tín chỉ"),
        ("Xếp loại", "fn_XepLoaiHocLuc → cập nhật SinhVien"),
        ("Hiển thị", "Client đọc DiemTB, HocLuc đã cập nhật"),
    ], size=14, numbered=True)
    code = """CREATE TRIGGER trg_DiemThi_TinhTongKet
ON DiemThi AFTER INSERT, UPDATE
AS
  UPDATE DiemThi SET TongKet =
    fn_TinhDiemTongKet(...)
  UPDATE SinhVien SET
    DiemTB = ..., HocLuc = ..."""
    _rect(s, Inches(6.85), Inches(1.3), Inches(5.85), Inches(3.2), CODE_BG, border=True, radius=True)
    _text(s, Inches(7.05), Inches(1.45), Inches(5.5), Inches(2.9), code, size=11, color=TEXT, font="Consolas")
    _image_placeholder(s, Inches(6.85), Inches(4.65), Inches(5.85), Inches(1.75), "frmChitietDiem.png")


def slide_client(prs):
    s = _blank(prs)
    _base_content(s, "Tien do Client", "Phan V")
    _bullets(s, Inches(0.65), Inches(1.3), Inches(5.8), Inches(5.2), [
        ("frmLogin", "Đăng nhập, xác thực tài khoản"),
        ("frmMain", "Menu điều hướng chính"),
        ("frmSv", "Danh sách, lọc lớp, tìm kiếm"),
        ("frmChitietSinhVien", "Hồ sơ, GPA, bảng điểm"),
        ("frmMonHoc", "CRUD môn học"),
        ("frmXepHang", "Báo cáo, thống kê xếp hạng"),
        ("Trạng thái", "~90% hoàn thiện giao diện"),
    ], size=14, numbered=True)
    _image_placeholder(s, Inches(6.7), Inches(1.25), Inches(6.0), Inches(5.3), "frmMain + frmSv (screenshot)")


def slide_server(prs):
    s = _blank(prs)
    _base_content(s, "Tien do Server", "Phan VI")
    _field_list(s, Inches(0.65), Inches(1.3), Inches(6.0), Inches(3.5), [
        ("Nền tảng", "ASP.NET Core Web API"),
        ("Hiện có", "Controllers, OpenAPI mẫu"),
        ("Kế hoạch", "Repository gọi Stored Procedure"),
        ("Hạ tầng", "Database Manager, middleware lỗi"),
        ("Bảo mật", "JWT / validation (giai đoạn sau)"),
    ])
    _rect(s, Inches(7.0), Inches(1.3), Inches(5.7), Inches(2.4), CODE_BG, border=True, radius=True)
    _text(s, Inches(7.2), Inches(1.45), Inches(5.3), Inches(2.1),
          "Client --HTTP--> Controller\n       |\n  Repository\n       |\n  EXEC SP --> SQL Server",
          size=13, color=TEXT, font="Consolas")
    _image_placeholder(s, Inches(7.0), Inches(3.9), Inches(5.7), Inches(2.6), "So do kien truc Server (ve them)")


def slide_database_tables(prs):
    s = _blank(prs)
    _base_content(s, "Co so du lieu - Bang chinh", "Phan VII")
    tables = [
        ("SinhVien", "MaSV, HoTen, MaLop, DiemTB, HocLuc"),
        ("MonHoc", "MaMH, TenMH, SoTinChi, GiangVien"),
        ("DiemThi", "MaSV, MaMH, QT, GK, CK, TongKet"),
        ("LopHoc", "MaLop, TenLop, Khoa"),
        ("TaiKhoan", "Dang nhap, QuyenHan"),
    ]
    y = Inches(1.35)
    for name, cols in tables:
        _rect(s, Inches(0.65), y, Inches(2.2), Inches(0.72), PRIMARY, radius=True)
        _text(s, Inches(0.65), y + Inches(0.14), Inches(2.2), Inches(0.45), name,
              size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_TITLE)
        _rect(s, Inches(2.95), y, Inches(9.5), Inches(0.72), WHITE, border=True, radius=True)
        _text(s, Inches(3.1), y + Inches(0.16), Inches(9.2), Inches(0.45), cols, size=13, color=TEXT)
        y += Inches(0.88)
    _text(s, Inches(0.65), Inches(6.0), Inches(12), Inches(0.4),
          "Du lieu mau: 20 SV, 15 mon, 5 lop  |  FK + CHECK diem 0-10", size=12, color=MUTED)


def slide_database_logic(prs):
    s = _blank(prs)
    _base_content(s, "Database - SP, Function, Trigger, View", "Phan VII")
    cols = [
        ("Procedure (9)", ["XacThucDangNhap", "LayDanhSachSinhVien", "TimKiemSinhVien",
                            "ThemSinhVien", "CapNhatDiem", "TinhDiemTrungBinh",
                            "LayBaoCaoXepHang", "DemSinhVienGioi", "LayDanhSachMonHoc"]),
        ("Function (2)", ["fn_TinhDiemTongKet", "fn_XepLoaiHocLuc"]),
        ("Trigger (1)", ["trg_DiemThi_TinhTongKet"]),
        ("View (2)", ["vw_BangXepHang", "vw_BaoCaoThanhTich"]),
    ]
    x = Inches(0.65)
    for title, items in cols:
        _rect(s, x, Inches(1.35), Inches(2.95), Inches(5.1), WHITE, border=True, radius=True)
        _rect(s, x, Inches(1.35), Inches(2.95), Inches(0.4), ACCENT, radius=True)
        _text(s, x, Inches(1.4), Inches(2.95), Inches(0.35), title, size=12, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER, font=FONT_TITLE)
        _bullets(s, x + Inches(0.15), Inches(1.85), Inches(2.65), Inches(4.4), items, size=11)
        x += Inches(3.15)


def slide_formula(prs):
    s = _blank(prs)
    _base_content(s, "Cong thuc tinh diem", "Phan VII")
    _rect(s, Inches(0.65), Inches(1.35), Inches(5.8), Inches(1.1), WHITE, border=True, radius=True)
    _text(s, Inches(0.85), Inches(1.55), Inches(5.4), Inches(0.75),
          "Diem mon = QT x 20% + GK x 30% + CK x 50%", size=18, bold=True, color=PRIMARY, font=FONT_TITLE)
    _rect(s, Inches(0.65), Inches(2.65), Inches(5.8), Inches(1.1), WHITE, border=True, radius=True)
    _text(s, Inches(0.85), Inches(2.85), Inches(5.4), Inches(0.75),
          "GPA = SUM(TongKet x TinChi) / SUM(TinChi)", size=16, bold=True, color=PRIMARY, font=FONT_TITLE)
    _bullets(s, Inches(0.65), Inches(4.0), Inches(5.8), Inches(2.2), [
        "Xuat sac >= 9.0",
        "Gioi >= 8.0",
        "Kha >= 6.5",
        "Trung binh >= 5.0",
        "Yeu < 5.0",
    ], size=14)
    _image_placeholder(s, Inches(6.7), Inches(1.35), Inches(6.0), Inches(5.0), "frmChitietDiem - cong thuc 20-30-50")


def slide_demo(prs):
    s = _blank(prs)
    _base_content(s, "Demo SQL Server", "Phan IX")
    _rect(s, Inches(0.65), Inches(1.3), Inches(6.2), Inches(4.8), CODE_BG, border=True, radius=True)
    sql = """EXEC LayBaoCaoXepHang
  @CheDo = 'TOP';

EXEC LayBaoCaoXepHang
  @CheDo = 'CANHBAO';

SELECT * FROM vw_BangXepHang;

EXEC CapNhatDiem
  @MaSV='SV001', @MaMH='IT03',
  @DiemQuaTrinh=8, @DiemGiuaKi=8,
  @DiemCuoiKi=9;"""
    _text(s, Inches(0.85), Inches(1.45), Inches(5.8), Inches(4.5), sql, size=12, color=TEXT, font="Consolas")
    _image_placeholder(s, Inches(7.1), Inches(1.3), Inches(5.6), Inches(4.8), "Screenshot SSMS ket qua query")


def slide_roadmap(prs):
    s = _blank(prs)
    _base_content(s, "Kế hoạch tiếp theo", "Phần IX")
    phases = [
        ("Tuần này", "Procedure Sửa/Xóa SV", "Nối Server: Login + danh sách SV", "Test trên máy nhóm"),
        ("Tuần sau", "Cập nhật điểm end-to-end", "Báo cáo xếp hạng trên UI", "Chuẩn hóa JSON API"),
    ]
    y = Inches(1.4)
    for phase, *tasks in phases:
        _rect(s, Inches(0.65), y, Inches(2.0), Inches(1.85), PRIMARY, radius=True)
        _text(s, Inches(0.65), y + Inches(0.65), Inches(2.0), Inches(0.5), phase,
              size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_TITLE)
        _rect(s, Inches(2.85), y, Inches(9.6), Inches(1.85), WHITE, border=True, radius=True)
        _bullets(s, Inches(3.05), y + Inches(0.2), Inches(9.2), Inches(1.5), list(tasks), size=14, numbered=True)
        y += Inches(2.05)


def slide_conclusion(prs):
    s = _blank(prs)
    _base_content(s, "Kết luận", "Phần IX")
    _bullets(s, Inches(0.65), Inches(1.35), Inches(6.0), Inches(4.5), [
        ("Client", "UI hoàn thiện, sẵn sàng tích hợp"),
        ("Database", "Logic đầy đủ: SP, trigger, view"),
        ("Server", "Cầu nối Client ↔ Database"),
        ("Trigger", "Tự động tính điểm, giảm lỗi nhập tay"),
        ("Tiếp theo", "Hoàn thiện pipeline end-to-end"),
    ], size=15, numbered=True)
    _image_placeholder(s, Inches(7.0), Inches(1.35), Inches(5.7), Inches(4.5), "Anh nhom / logo du an")


def slide_closing(prs):
    s = _blank(prs)
    _rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, PRIMARY)
    _rect(s, Inches(2.2), Inches(2.5), Inches(8.9), Inches(2.5), PRIMARY_SOFT, border=True, radius=True)
    _text(s, Inches(2.4), Inches(2.85), Inches(8.5), Inches(0.9),
          "CẢM ƠN THẦY VÀ CÁC BẠN!", size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_TITLE)
    _text(s, Inches(2.4), Inches(3.75), Inches(8.5), Inches(0.55),
          "Hỏi & đáp", size=20, color=ACCENT, align=PP_ALIGN.CENTER)


def main():
    global _footer_total, _footer_idx
    _footer_idx = 0
    # Slides co footer (khong tinh cover, section, closing)
    _footer_total = 16

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_agenda(prs)

    slide_section(prs, "I", "Giới thiệu", "Đề tài và mục tiêu")
    slide_intro(prs)

    slide_section(prs, "II", "Cấu trúc dự án", "Thư mục và phân chia công việc")
    slide_folder_tree(prs)

    slide_section(prs, "III", "Kiến trúc & Pipeline", "Luồng xử lý yêu cầu")
    slide_pipeline(prs)

    slide_section(prs, "IV", "Chức năng ứng dụng", "Ứng dụng làm gì?")
    slide_app_features(prs)

    slide_section(prs, "V", "Tiến độ Client", "Giao diện WinForms")
    slide_client(prs)

    slide_section(prs, "VI", "Tiến độ Server", "ASP.NET Core API")
    slide_server(prs)

    slide_section(prs, "VII", "Cơ sở dữ liệu", "Schema và logic SQL")
    slide_database_tables(prs)
    slide_database_logic(prs)
    slide_formula(prs)

    slide_section(prs, "VIII", "Trigger tự động", "Cập nhật điểm thời gian thực")
    slide_trigger_flow(prs)

    slide_section(prs, "IX", "Demo & Kế hoạch", "Trình bày và hướng phát triển")
    slide_demo(prs)
    slide_roadmap(prs)
    slide_conclusion(prs)

    slide_closing(prs)

    out = OUTPUT
    try:
        prs.save(out)
    except PermissionError:
        out = OUTPUT.with_name(OUTPUT.stem + "_v2" + OUTPUT.suffix)
        prs.save(out)
        print("Note: Dong file PPTX cu neu dang mo, roi chay lai de ghi de.")
    print(f"Created: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
